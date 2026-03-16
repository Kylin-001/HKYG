#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 165, 0)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(12)
        p.space_after = Pt(6)
    
    return slide

def add_info_slide(prs, title, sections):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    y_pos = 1.3
    for section_title, items in sections:
        section_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12.333), Inches(0.5))
        tf = section_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 165, 0)
        y_pos += 0.5
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12.333), Inches(0.8))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.space_before = Pt(6)
        y_pos += len(items) * 0.35 + 0.3
    
    return slide

def add_architecture_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "三、系统架构设计"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    layers = [
        ("表现层", "微信小程序 / Web应用 / 管理后台"),
        ("API网关层", "Spring Cloud Gateway"),
        ("业务服务层", "11个微服务（用户/商品/订单/支付等）"),
        ("数据访问层", "MySQL 8.0 / Redis 4.4 / MyBatis Plus"),
        ("基础设施层", "Nacos / 消息队列 / Sentry监控")
    ]
    
    y_pos = 1.5
    for layer_name, layer_desc in layers:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(y_pos),
            Inches(9.333), Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
        shape.line.color.rgb = RGBColor(0, 112, 192)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{layer_name}\n{layer_desc}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if y_pos < 5.5:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(6.166), Inches(y_pos + 0.85),
                Inches(0.5), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(255, 165, 0)
            arrow.line.color.rgb = RGBColor(255, 165, 0)
        
        y_pos += 1.2
    
    return slide

def add_module_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "四、核心功能模块"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    modules = [
        ("用户中心", "注册登录\n实名认证\n权限控制"),
        ("商品商城", "商品浏览\n购物车下单\n智能推荐"),
        ("校园跑腿", "取快递\n代购代办\n悬赏接单"),
        ("外卖服务", "外卖柜配送\n寝室配送\n订单管理"),
        ("二手交易", "商品发布\n在线议价\n交易保障"),
        ("失物招领", "失物发布\n寻物启事\n消息通知"),
        ("营销系统", "优惠券\n积分系统\n会员等级"),
        ("后台管理", "数据看板\n用户管理\n系统配置")
    ]
    
    x_positions = [0.5, 3.5, 6.5, 9.5]
    y_positions = [1.5, 4.2]
    
    for i, (module_name, module_desc) in enumerate(modules):
        x = x_positions[i % 4]
        y = y_positions[i // 4]
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(2.8), Inches(2.3)
        )
        shape.fill.solid()
        if i < 4:
            shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
        else:
            shape.fill.fore_color.rgb = RGBColor(255, 165, 0)
        shape.line.color.rgb = RGBColor(0, 0, 0)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = module_name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = module_desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_role_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "五、用户角色与用例图"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    roles = [
        ("学生用户", ["注册登录、实名认证", "浏览商品、购物下单", "发布跑腿需求、接单配送", "二手交易、失物招领", "查看订单、评价互动"]),
        ("商家用户", ["商家入驻申请", "商品信息管理", "订单处理与发货", "营业数据分析", "营销活动设置"]),
        ("配送员用户", ["配送员注册审核", "查看附近跑腿订单", "接单配送、状态更新", "收入管理与提现", "服务评价查看"]),
        ("管理员用户", ["数据看板与统计分析", "用户管理与审核", "商品与商家审核", "订单管理与纠纷处理", "内容发布与系统配置"])
    ]
    
    x_positions = [0.5, 6.5]
    y_positions = [1.3, 4.2]
    
    for i, (role_name, role_items) in enumerate(roles):
        x = x_positions[i % 2]
        y = y_positions[i // 2]
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(6), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
        shape.line.color.rgb = RGBColor(0, 112, 192)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = role_name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        for item in role_items:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(4)
    
    return slide

def add_tech_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "六、技术栈"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    backend_techs = [
        "Spring Boot 3.2.2", "Spring Cloud 2023.0.0", "Nacos 2.0.0",
        "Spring Cloud Gateway 4.1.1", "MyBatis Plus 3.5.5", "Spring Security + JWT",
        "MySQL 8.0.33", "Redis 4.4.3", "SpringDoc OpenAPI 3"
    ]
    
    frontend_techs = [
        "Vue 3.5.1", "Vite 5.0.0", "Element Plus 2.4.4",
        "Pinia 2.1.7", "Vue Router 4.2.5", "Axios 1.6.0",
        "ECharts 5.4.0", "TypeScript 5.2.2", "Vue I18n 12.0.0"
    ]
    
    backend_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.3),
        Inches(6), Inches(3.5)
    )
    backend_shape.fill.solid()
    backend_shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
    backend_shape.line.color.rgb = RGBColor(0, 112, 192)
    
    tf = backend_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "后端技术栈"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    for tech in backend_techs:
        p = tf.add_paragraph()
        p.text = "• " + tech
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(4)
    
    frontend_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.833), Inches(1.3),
        Inches(6), Inches(3.5)
    )
    frontend_shape.fill.solid()
    frontend_shape.fill.fore_color.rgb = RGBColor(255, 165, 0)
    frontend_shape.line.color.rgb = RGBColor(255, 165, 0)
    
    tf = frontend_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "前端技术栈"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    for tech in frontend_techs:
        p = tf.add_paragraph()
        p.text = "• " + tech
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(4)
    
    features_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(2))
    tf = features_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "核心特性"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 165, 0)
    
    features = [
        "前后端分离架构 - 支持Web端、管理后台、微信小程序多端访问",
        "微服务架构 - 11个业务服务独立部署，高内聚低耦合",
        "协同过滤推荐算法 - 个性化商品推荐，提升用户体验"
    ]
    for feature in features:
        p = tf.add_paragraph()
        p.text = "• " + feature
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(6)
    
    return slide

def add_database_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "七、数据库设计"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    db_modules = [
        ("用户模块", "用户表、用户认证表、角色表、权限表、地址表", "用户信息、权限管理、地址管理"),
        ("商品模块", "商品分类表、商品表、购物车表", "商品管理、购物车"),
        ("订单模块", "订单主表、订单明细表、支付记录表", "订单处理、支付记录"),
        ("配送模块", "跑腿需求表、外卖订单表、配送员表、外卖柜表", "配送管理、外卖柜管理"),
        ("营销模块", "优惠券表、用户优惠券表、积分记录表、会员等级表", "营销活动、会员管理"),
        ("其他模块", "失物表、招领表、二手商品表、校园信息表", "失物招领、二手交易、校园服务")
    ]
    
    y_pos = 1.3
    for module_name, tables, desc in db_modules:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Inches(12.333), Inches(0.9)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
        shape.line.color.rgb = RGBColor(0, 112, 192)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{module_name}：{tables}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        p = tf.add_paragraph()
        p.text = f"功能：{desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        y_pos += 1.0
    
    return slide

def add_feasibility_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "八、可行性分析"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    feasibilities = [
        ("技术可行性 ✓", ["采用主流技术栈，文档完善", "核心功能已完成100%", "11个微服务已独立部署", "基础设施搭建完成"]),
        ("经济可行性 ✓", ["开源技术，无版权费用", "硬件要求低，成本可控", "盈利模式清晰（佣金/广告）", "可商业化适配转化"]),
        ("操作可行性 ✓", ["界面简洁直观，易于上手", "符合主流互联网产品习惯", "支持多端访问", "微信小程序无需安装"]),
        ("法律可行性 ✓", ["遵守网络安全法等法规", "用户信息加密保护", "完善的用户协议与隐私政策", "实名认证保障交易安全"])
    ]
    
    x_positions = [0.5, 6.5]
    y_positions = [1.3, 4.2]
    
    for i, (feasibility_name, items) in enumerate(feasibilities):
        x = x_positions[i % 2]
        y = y_positions[i // 2]
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(6), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
        shape.line.color.rgb = RGBColor(0, 112, 192)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = feasibility_name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        for item in items:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(4)
    
    return slide

def add_timeline_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "九、时间规划"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    timeline = [
        ("第1-2周", "文献调研\n需求分析"),
        ("第3-4周", "系统设计\n环境搭建"),
        ("第5-8周", "后端开发\n单元测试"),
        ("第9-12周", "前端开发\n前后端联调"),
        ("第13-14周", "算法优化\n系统测试"),
        ("第15-16周", "部署上线\n文档撰写"),
        ("第17周", "答辩准备\n毕业答辩"),
        ("预期成果", "完整系统\n论文文档")
    ]
    
    x_positions = [0.5, 3.5, 6.5, 9.5]
    y_positions = [1.3, 3.8]
    
    for i, (week, desc) in enumerate(timeline):
        x = x_positions[i % 4]
        y = y_positions[i // 4]
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(2.8), Inches(2)
        )
        shape.fill.solid()
        if i < 7:
            shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
        else:
            shape.fill.fore_color.rgb = RGBColor(255, 165, 0)
        shape.line.color.rgb = RGBColor(0, 0, 0)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = week
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    results_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12.333), Inches(1.2))
    tf = results_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "预期成果"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 165, 0)
    
    results = [
        "代码规模：后端约50,000行，前端约30,000行",
        "系统规模：Java类约300个，Vue组件约200个，API接口约200个",
        "测试覆盖：测试用例约250个，覆盖率约77%"
    ]
    for result in results:
        p = tf.add_paragraph()
        p.text = "• " + result
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(4)
    
    return slide

def add_final_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢聆听"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "敬请各位老师批评指正"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 165, 0)
    p.alignment = PP_ALIGN.CENTER
    
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "答辩人：张开源"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "指导教师：李懿 讲师"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

slide1 = add_title_slide(prs, 
    "基于Java的黑科易购校园服务平台\n设计与实现",
    "本科毕业设计开题答辩"
)

info_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(2))
tf = info_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "答辩人：张开源    学号：2022024874"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0, 0, 0)
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "班级：数据22-4班    学院：计算机与信息工程学院（软件学院）"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0, 0, 0)
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "指导教师：李懿 讲师"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0, 0, 0)
p.alignment = PP_ALIGN.CENTER

add_info_slide(prs, "一、研究背景与意义", [
    ("研究背景", [
        "大三、大四学生时间紧张，校园服务分散",
        "现有平台功能单一，服务质量参差不齐",
        "校园资源缺乏有效整合与利用"
    ]),
    ("研究意义", [
        "社会意义：整合校园服务资源，提升学生生活效率",
        "技术意义：实践微服务架构，探索校园场景创新"
    ])
])

add_info_slide(prs, "二、研究目标", [
    ("总体目标", [
        "设计并实现面向黑龙江科技大学师生的综合性校园服务平台",
        "整合购物、外卖、跑腿、二手交易、失物招领等功能"
    ]),
    ("功能目标", [
        "多角色用户体系（学生/商家/配送员/管理员）",
        "智能推荐系统（协同过滤算法）",
        "营销系统（优惠券/积分/会员）"
    ]),
    ("性能目标", [
        "API响应时间 < 200ms",
        "支持1000+并发用户",
        "系统可用性 > 99.9%"
    ])
])

add_architecture_slide(prs)
add_module_slide(prs)
add_role_slide(prs)
add_tech_slide(prs)
add_database_slide(prs)
add_feasibility_slide(prs)
add_timeline_slide(prs)
add_final_slide(prs)

output_path = '/home/zky/HKYG/开题答辩PPT.pptx'
prs.save(output_path)
print(f"PPT已成功生成：{output_path}")
