#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
基于项目实际完成情况生成更新版开题答辩PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def create_updated_opening_report_ppt():
    """创建基于项目实际完成情况的开题答辩PPT"""
    
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 定义颜色主题
    primary_color = RGBColor(0x19, 0x76, 0xD2)  # 蓝色
    secondary_color = RGBColor(0xFF, 0x98, 0x00)  # 橙色
    accent_color = RGBColor(0x4C, 0xAF, 0x50)  # 绿色
    
    # ===== 第1页：封面 =====
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "黑科易购——黑龙江科技大学特色校园服务平台"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "的设计与实现"
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.alignment = PP_ALIGN.CENTER
    
    # 答辩信息
    info_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(1.5))
    tf = info_box.text_frame
    tf.word_wrap = True
    
    infos = [
        "2026届本科毕业设计（论文）开题答辩",
        "",
        "学生姓名：张开源",
        "学    号：2022024874",
        "班    级：数据22-4班",
        "指导教师：李懿 讲师",
        "答辩日期：2025年3月"
    ]
    
    for i, info in enumerate(infos):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = info
        p.font.size = Pt(20) if i == 0 else Pt(18)
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p.alignment = PP_ALIGN.CENTER
    
    # ===== 第2页：目录 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
    slide.shapes.title.text = "答辩提纲"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    items = [
        "一、选题背景与意义",
        "二、国内外研究现状",
        "三、研究目标与内容",
        "四、技术方案与架构",
        "五、功能模块设计",
        "六、项目完成情况",
        "七、预期成果与创新点"
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = content.paragraphs[0]
        else:
            p = content.add_paragraph()
        p.text = item
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(12)
    
    # ===== 第3页：选题背景 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "一、选题背景与意义"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    bg_text = [
        "研究背景：",
        "• 移动互联网普及，高校师生对校园服务便捷性要求提高",
        "• 现有校园服务入口分散、信息孤岛、用户体验差",
        "• 缺乏针对校园场景的特色服务功能",
        "",
        "研究意义：",
        "• 理论意义：探索微服务架构在校园信息化中的应用",
        "• 实践意义：为师生提供一站式校园服务平台",
        "• 社会意义：整合校园服务资源，提升学生生活效率"
    ]
    
    for i, text in enumerate(bg_text):
        if i == 0:
            p = content.paragraphs[0]
        else:
            p = content.add_paragraph()
        p.text = text
        p.font.size = Pt(24) if text.endswith('：') else Pt(20)
        p.font.bold = text.endswith('：')
        p.font.color.rgb = primary_color if text.endswith('：') else RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(8)
    
    # ===== 第4页：国内外研究现状 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "二、国内外研究现状"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    research_text = [
        "国外研究：",
        "• MIT、斯坦福等高校采用一体化校园管理系统",
        "• 剑桥大学、牛津大学开发完善的校园移动应用",
        "• 技术特点：云原生架构、注重用户体验和数据安全",
        "",
        "国内研究：",
        "• 清华大学：开发\"清华校园\"APP，整合多方位服务",
        "• 浙江大学：推出\"浙大钉\"，实现移动端一站式办公",
        "• 商业平台：美团、饿了么提供校园外卖，但缺乏定制化",
        "",
        "发展趋势：微服务架构、多端融合、智能化、数据驱动"
    ]
    
    for i, text in enumerate(research_text):
        if i == 0:
            p = content.paragraphs[0]
        else:
            p = content.add_paragraph()
        p.text = text
        p.font.size = Pt(22) if text.endswith('：') else Pt(18)
        p.font.bold = text.endswith('：')
        p.font.color.rgb = primary_color if text.endswith('：') else RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(6)
    
    # ===== 第5页：研究目标 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "三、研究目标与内容"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    goal_text = [
        "研究目标：",
        "设计并实现基于微服务架构的黑龙江科技大学特色校园服务平台",
        "实现校园外卖、商品购物、二手交易、失物招领等核心功能",
        "",
        "研究内容：",
        "1. 需求分析：调研师生校园服务需求",
        "2. 架构设计：设计高可用、可扩展的微服务架构",
        "3. 功能实现：用户、商品、订单、支付、外卖等核心模块",
        "4. 性能优化：研究高并发场景下的性能优化策略",
        "5. 安全保障：设计完善的安全机制"
    ]
    
    for i, text in enumerate(goal_text):
        if i == 0:
            p = content.paragraphs[0]
        else:
            p = content.add_paragraph()
        p.text = text
        p.font.size = Pt(24) if text.endswith('：') else Pt(20)
        p.font.bold = text.endswith('：')
        p.font.color.rgb = primary_color if text.endswith('：') else RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(8)
    
    # ===== 第6页：技术方案 - 系统架构 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "四、技术方案 - 系统架构"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    arch_text = [
        "五层微服务架构：",
        "",
        "① 前端层：Web管理端(Vue3) + 移动端/小程序",
        "② API网关层：Spring Cloud Gateway（路由/限流/认证）",
        "③ 服务层：11个微服务（用户/商品/订单/支付/外卖等）",
        "④ 数据层：MySQL集群 + Redis缓存 + RabbitMQ消息队列",
        "⑤ 基础设施层：Docker + Nacos + Sentry监控"
    ]
    
    for i, text in enumerate(arch_text):
        if i == 0:
            p = content.paragraphs[0]
        else:
            p = content.add_paragraph()
        p.text = text
        p.font.size = Pt(26) if i == 0 else Pt(20)
        p.font.bold = i == 0
        p.font.color.rgb = primary_color if i == 0 else RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(10)
    
    # ===== 第7页：技术方案 - 技术栈 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "四、技术方案 - 核心技术栈"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    tech_text = [
        "后端技术：",
        "• Java 17 + Spring Boot 3.2.2 + Spring Cloud 2023.0.0",
        "• Spring Cloud Alibaba 2023.0.1.0 (Nacos、Gateway、OpenFeign)",
        "• MyBatis-Plus 3.5.5 + MySQL 8.0.33 + Redis 4.4.3",
        "• Spring Security + JWT + SpringDoc OpenAPI 3",
        "",
        "前端技术：",
        "• Vue.js 3.5.1 + TypeScript 5.2.2 + Vite 5.0.0",
        "• Element Plus 2.4.4 (PC端) + 微信小程序",
        "• Pinia 2.1.7 + Vue Router 4.2.5 + Axios 1.6.0",
        "",
        "运维技术：Docker + Nacos + Sentry监控"
    ]
    
    for i, text in enumerate(tech_text):
        if i == 0:
            p = content.paragraphs[0]
        else:
            p = content.add_paragraph()
        p.text = text
        p.font.size = Pt(22) if text.endswith('：') else Pt(18)
        p.font.bold = text.endswith('：')
        p.font.color.rgb = secondary_color if text.endswith('：') else RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(6)
    
    # ===== 第8页：功能模块设计 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "五、功能模块设计"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    module_text = [
        "用户端：",
        "• 微信小程序：商品浏览、下单支付、订单跟踪、个人中心",
        "• 二手市场、失物招领、智能客服",
        "",
        "管理端：",
        "• Web管理后台：用户管理、商品审核、订单监控、数据统计",
        "• 营销活动管理、系统配置",
        "",
        "后端微服务（11个）：",
        "• 用户服务、商品服务、订单服务、支付服务、营销服务",
        "• 外卖服务、配送服务、校园服务、二手服务、失物招领服务"
    ]
    
    for i, text in enumerate(module_text):
        if i == 0:
            p = content.paragraphs[0]
        else:
            p = content.add_paragraph()
        p.text = text
        p.font.size = Pt(22) if text.endswith('：') else Pt(18)
        p.font.bold = text.endswith('：')
        p.font.color.rgb = accent_color if text.endswith('：') else RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(8)
    
    # ===== 第9页：项目完成情况（新增） =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "六、项目完成情况"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    progress_text = [
        "项目整体进度：✅ 已完成 100%",
        "",
        "核心功能完成情况：",
        "• P0任务（高优先级）：Vue 3迁移、测试体系、性能优化、安全加固 ✅",
        "• 用户认证优化、权限控制、商品推荐系统、支付安全增强 ✅",
        "• 营销系统开发、数据分析系统、协同过滤算法、系统监控 ✅",
        "",
        "代码统计：",
        "• 后端代码：约50,000行（Java类约300个）",
        "• 前端代码：约30,000行（Vue组件约200个）",
        "• API接口：约200个 | 数据库表：约30个 | 测试用例：约250个"
    ]
    
    for i, text in enumerate(progress_text):
        if i == 0:
            p = content.paragraphs[0]
        else:
            p = content.add_paragraph()
        p.text = text
        p.font.size = Pt(24) if i == 0 else Pt(18)
        p.font.bold = i == 0
        p.font.color.rgb = accent_color if i == 0 else RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(6)
    
    # ===== 第10页：预期成果与创新点 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "七、预期成果与创新点"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    result_text = [
        "预期成果：",
        "• 软件系统：Web管理端 + 微信小程序 + 后端微服务（11个）",
        "• 技术文档：需求规格、设计文档、接口文档、测试报告",
        "• 毕业论文：一篇高质量的本科毕业论文",
        "",
        "创新点：",
        "1. 校园特色功能：食堂外卖、校园地图、学业辅助等定制化功能",
        "2. 微服务架构：Spring Cloud Alibaba实现高可用、可扩展",
        "3. 智能推荐：基于协同过滤算法的商品智能推荐系统",
        "4. 多端统一：Vue 3 + 微信小程序实现多端访问",
        "5. 实时配送：集成地图服务，实现配送员实时位置追踪"
    ]
    
    for i, text in enumerate(result_text):
        if i == 0:
            p = content.paragraphs[0]
        else:
            p = content.add_paragraph()
        p.text = text
        p.font.size = Pt(22) if text.endswith('：') else Pt(18)
        p.font.bold = text.endswith('：')
        p.font.color.rgb = primary_color if text.endswith('：') else RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(6)
    
    # ===== 第11页：性能指标 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "八、系统性能指标"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    perf_text = [
        "响应时间：",
        "• API平均响应时间：< 200ms",
        "• 页面加载时间：< 2s",
        "• 数据库查询时间：< 100ms",
        "",
        "并发性能：",
        "• 支持并发用户数：> 1000",
        "• QPS：> 500",
        "• 系统可用性：> 99.9%",
        "",
        "安全指标：",
        "• 密码加密：BCrypt | 传输加密：HTTPS",
        "• 认证方式：JWT | 权限控制：RBAC",
        "• Redis缓存命中率：> 90%"
    ]
    
    for i, text in enumerate(perf_text):
        if i == 0:
            p = content.paragraphs[0]
        else:
            p = content.add_paragraph()
        p.text = text
        p.font.size = Pt(22) if text.endswith('：') else Pt(18)
        p.font.bold = text.endswith('：')
        p.font.color.rgb = secondary_color if text.endswith('：') else RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(6)
    
    # ===== 第12页：参考文献 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "主要参考文献"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    ref_text = [
        "[1] Turnquist G L, et al. Microservices with Spring Boot 3 and Spring Cloud[M]. 3rd ed. Packt, 2022.",
        "[2] Meric A. Mastering Spring Boot 3.0[M]. Packt Publishing, 2024.",
        "[3] Raj P, et al. Microservices Design Patterns with Java[M]. Packt, 2024.",
        "[4] Sharma S. Modern API Development with Spring 6 and Spring Boot 3[M]. Packt, 2023.",
        "[5] Hinkula J. Full Stack Development with Spring Boot 3 and React[M]. Packt, 2023."
    ]
    
    for i, text in enumerate(ref_text):
        if i == 0:
            p = content.paragraphs[0]
        else:
            p = content.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(8)
    
    # ===== 第13页：致谢 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(2))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢各位老师的聆听与指导！"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p.alignment = PP_ALIGN.CENTER
    
    # 保存PPT
    output_path = "开题答辩PPT_2025更新版.pptx"
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")
    print(f"总页数: {len(prs.slides)}")
    print(f"幻灯片内容概览：")
    for i, slide in enumerate(prs.slides, 1):
        try:
            title = slide.shapes.title.text if slide.shapes.title else f"第{i}页"
            print(f"  第{i}页: {title}")
        except:
            print(f"  第{i}页: （致谢页）")
    
    return output_path

if __name__ == "__main__":
    create_updated_opening_report_ppt()
